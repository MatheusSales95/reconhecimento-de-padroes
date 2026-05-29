"""
Gera RF_modelo.pptx no mesmo estilo visual do REC PADRÕES.
Pipeline: Título → Como Funciona → Entradas → Desbalanceamento
          → Validação Espacial → Estratégias → Threshold
          → Saída (Mapa) → Saída (Erros) → Importância
"""
import os, warnings
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
warnings.filterwarnings('ignore')

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from PIL import Image as PILImage

# ── Caminhos ────────────────────────────────────────────────────────────────
BASE    = '/home/matheus-sales/Documents/atividade_reconhecimento'
RF_FIGS = f'{BASE}/relatorio/figs/rf'
GEN     = f'{BASE}/relatorio/pptx_figs'
ORIG    = f'{BASE}/relatorio/Classificação de áreas queimadas - REC PADRÕES.pptx'
OUT     = f'{BASE}/relatorio/RF_modelo.pptx'

P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'

# ── Cores do tema pedro/REC ──────────────────────────────────────────────────
BLUE  = RGBColor(0x42, 0x85, 0xF4)
DARK  = RGBColor(0x21, 0x21, 0x21)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY  = RGBColor(0x55, 0x55, 0x55)

# ── Dimensões do slide (cm) ──────────────────────────────────────────────────
SW, SH = 25.40, 14.29

# ============================================================
# PARTE 1 — FIGURAS GERADAS (tabelas como imagem)
# ============================================================
os.makedirs(GEN, exist_ok=True)

def savefig(name, dpi=150):
    path = f'{GEN}/{name}.png'
    plt.savefig(path, dpi=dpi, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close('all')
    print(f'  ✓ {name}.png')
    return path

print('=== Gerando tabelas como imagem ===')

# ── Tabela: Estratégias de desbalanceamento ──────────────────
fig, ax = plt.subplots(figsize=(11, 2.8))
ax.axis('off')
headers = ['Estratégia', 'F₁', 'MCC', 'Kappa', 'Prec.', 'Recall', 'AUPRC']
rows = [
    ['Sem tratamento',          '0,033', '0,088', '0,030', '0,535', '0,017', '0,202'],
    ['class_weight (τ=0.5)',    '0,026', '0,086', '0,024', '0,632', '0,013', '0,191'],
    ['SMOTE',                   '0,220', '0,173', '0,172', '0,200', '0,245', '0,154'],
    ['SMOTETomek',              '0,220', '0,173', '0,172', '0,200', '0,245', '0,158'],
    ['Undersampling',           '0,201', '0,198', '0,122', '0,117', '0,716', '0,216'],
    ['class_weight + τ*',       '0,262', '0,217', '0,214', '0,231', '0,302', '0,191'],
]
colors = [['#f0f4ff']*7 if i%2==0 else ['#ffffff']*7 for i in range(len(rows))]
colors[-1] = ['#d4edda']*7   # destaque linha vencedora
tbl = ax.table(cellText=rows, colLabels=headers, cellLoc='center',
               loc='center', cellColours=colors)
tbl.auto_set_font_size(False)
tbl.set_fontsize(10)
tbl.scale(1, 1.7)
for j in range(len(headers)):
    tbl[0, j].set_facecolor('#4285F4')
    tbl[0, j].set_text_props(color='white', fontweight='bold')
tbl[len(rows), 1].set_text_props(fontweight='bold')   # best F1
plt.tight_layout()
TAB_ESTRATEGIAS = savefig('tab_estrategias_rf')

# ── Tabela: RF otimizado (threshold) ────────────────────────
fig, ax = plt.subplots(figsize=(10, 1.8))
ax.axis('off')
headers2 = ['Configuração', 'F₁', 'MCC', 'Kappa', 'Prec.', 'Recall', 'τ']
rows2 = [
    ['RF Otim. τ = 0,5',  '0,301', '0,267', '0,248', '0,224', '0,459', '0,50'],
    ['RF Otim. τ* = 0,65','0,329', '0,299', '0,297', '0,378', '0,291', '0,65'],
]
c2 = [['#f0f4ff']*7, ['#d4edda']*7]
tbl2 = ax.table(cellText=rows2, colLabels=headers2, cellLoc='center',
                loc='center', cellColours=c2)
tbl2.auto_set_font_size(False)
tbl2.set_fontsize(11)
tbl2.scale(1, 2.2)
for j in range(len(headers2)):
    tbl2[0, j].set_facecolor('#4285F4')
    tbl2[0, j].set_text_props(color='white', fontweight='bold')
tbl2[2, 1].set_text_props(fontweight='bold')
plt.tight_layout()
TAB_THRESHOLD = savefig('tab_threshold_rf')

# ── Tabela: Validação cross-cena ─────────────────────────────
fig, ax = plt.subplots(figsize=(10, 1.8))
ax.axis('off')
headers3 = ['Treino → Teste', 'F₁', 'MCC', 'Kappa', 'Prec.', 'Recall', 'AUPRC']
rows3 = [
    ['221/074 → Intra (blocos)', '0,329', '0,299', '0,297', '0,378', '0,291', '0,260'],
    ['221/074 → 220/075 (cross)','0,411', '0,374', '0,369', '0,486', '0,356', '0,467'],
]
c3 = [['#f0f4ff']*7, ['#d4edda']*7]
tbl3 = ax.table(cellText=rows3, colLabels=headers3, cellLoc='center',
                loc='center', cellColours=c3)
tbl3.auto_set_font_size(False)
tbl3.set_fontsize(11)
tbl3.scale(1, 2.2)
for j in range(len(headers3)):
    tbl3[0, j].set_facecolor('#4285F4')
    tbl3[0, j].set_text_props(color='white', fontweight='bold')
tbl3[2, 1].set_text_props(fontweight='bold')
plt.tight_layout()
TAB_CROSSCENA = savefig('tab_crosscena_rf')

# ── Tabela: Mapa de predição completo ───────────────────────
fig, ax = plt.subplots(figsize=(7, 2.0))
ax.axis('off')
headers4 = ['', 'Ref: Queimado', 'Ref: Não queimado']
rows4 = [
    ['Pred: Queimado',     'TP = 16.252', 'FP = 39.328'],
    ['Pred: Não queimado', 'FN = 59.905', 'TN = 1.470.429'],
    ['Precision',  '29,24%', ''],
    ['Recall',     '21,34%', ''],
    ['F₁',         '0,247',  ''],
]
c4 = [['#fff3cd','#fff3cd','#fff3cd'],
      ['#fff3cd','#fff3cd','#fff3cd'],
      ['#f0f4ff','#f0f4ff','#f0f4ff'],
      ['#f0f4ff','#f0f4ff','#f0f4ff'],
      ['#d4edda','#d4edda','#d4edda']]
tbl4 = ax.table(cellText=rows4, colLabels=headers4, cellLoc='center',
                loc='center', cellColours=c4)
tbl4.auto_set_font_size(False)
tbl4.set_fontsize(10)
tbl4.scale(1, 1.8)
for j in range(3):
    tbl4[0, j].set_facecolor('#4285F4')
    tbl4[0, j].set_text_props(color='white', fontweight='bold')
plt.tight_layout()
TAB_MAPA = savefig('tab_mapa_rf')

# ── Ilustração: entradas do modelo ──────────────────────────
# Geometria (tudo em axes fraction 0-1):
#   margin=0.01, bw=0.09 (half-width), gap=0.07 entre caixas
#   c1=0.10, c2=0.35, c3=0.60, c4=0.85
#   setas nos gaps: [0.20→0.25], [0.45→0.50], [0.70→0.75]
fig, ax = plt.subplots(figsize=(13, 3.8))
ax.axis('off')
ax.set_xlim(0, 1); ax.set_ylim(0, 1)

BW   = 0.09   # half-width de cada caixa
Y    = 0.60   # linha central
YTOP = 0.20   # borda superior da caixa
YBOT = 0.82   # borda inferior da caixa

BOX_DATA = [
    (0.10, 'Landsat 8/9\nPré-fogo  /  Pós-fogo', '#4285F4'),
    (0.35, 'Extração de\nAtributos\n10 features/pixel', '#34A853'),
    (0.60, 'Random\nForest\n171 árvores', '#EA4335'),
    (0.85, 'Mapa de\nQueimadas\n(binário)', '#FBBC05'),
]

# 1) Caixas (zorder baixo — ficam atrás)
for xc, txt, col in BOX_DATA:
    rect = plt.Rectangle((xc - BW, YTOP), 2 * BW, YBOT - YTOP,
                          facecolor=col, alpha=0.13, edgecolor=col, lw=2,
                          transform=ax.transAxes, clip_on=False, zorder=1)
    ax.add_patch(rect)
    ax.text(xc, Y, txt, ha='center', va='center', fontsize=11,
            fontweight='bold', color=col, transform=ax.transAxes, zorder=2)

# 2) Setas nos gaps entre caixas (zorder alto — ficam na frente)
arrow_kw = dict(arrowstyle='->', color='#333', lw=2.5, mutation_scale=20)
for x0, x1 in [(0.20, 0.25), (0.45, 0.50), (0.70, 0.75)]:
    ax.annotate('', xy=(x1, Y), xytext=(x0, Y),
                xycoords='axes fraction', textcoords='axes fraction',
                arrowprops=arrow_kw, zorder=4)

# 3) Legenda de features
feat_txt = ('Atributos: B4 · B5 · B7 · NDVI · NBR  '
            '(pré-fogo + pós-fogo = 10 atributos por pixel)   ·   ~39 milhões de pixels')
ax.text(0.5, 0.03, feat_txt, ha='center', va='bottom', fontsize=9,
        color='#555', transform=ax.transAxes, style='italic', zorder=2)

plt.tight_layout(pad=0.4)
PIPELINE_IMG = savefig('pipeline_entrada_rf')

print('✓ Tabelas geradas\n')

# ============================================================
# PARTE 2 — CONSTRUIR O PPTX
# ============================================================
print('=== Construindo PPTX ===')


def clear_slide(slide):
    sp_tree = slide.shapes._spTree
    for child in list(sp_tree):
        tag = etree.QName(child.tag).localname
        if tag == 'pic':
            sp_tree.remove(child)
        elif tag == 'sp':
            nvPr = child.find(f'.//{{{P_NS}}}nvPr')
            ph   = nvPr.find(f'{{{P_NS}}}ph') if nvPr is not None else None
            if ph is None:
                sp_tree.remove(child)
        elif tag in ('grpSp', 'graphicFrame', 'cxnSp'):
            sp_tree.remove(child)
    for ph in slide.placeholders:
        if ph.has_text_frame:
            ph.text_frame.clear()


def add_title(slide, text, fontsize=20):
    """Título pequeno no canto superior esquerdo (estilo REC PADRÕES)."""
    txBox = slide.shapes.add_textbox(Cm(0.87), Cm(0.18), Cm(15), Cm(1.09))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fontsize)
    run.font.bold = True
    run.font.color.rgb = BLACK


def add_center_title(slide, title, subtitle='', authors=''):
    """Slide de título centralizado (estilo slide 1 do REC)."""
    # Título central
    txBox = slide.shapes.add_textbox(Cm(4.5), Cm(4.5), Cm(16.4), Cm(4.05))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = title
    run.font.size = Pt(28); run.font.bold = True; run.font.color.rgb = BLACK
    # Subtítulo/authors no rodapé
    if subtitle or authors:
        txBox2 = slide.shapes.add_textbox(Cm(0.87), Cm(9.8), Cm(23.67), Cm(4.0))
        tf2 = txBox2.text_frame; tf2.word_wrap = True
        for line in [subtitle, authors]:
            if not line: continue
            p2 = tf2.add_paragraph() if tf2.text else tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.add_run(); run2.text = line
            run2.font.size = Pt(16); run2.font.color.rgb = GREY


def add_section_title(slide, text):
    """Slide de seção com texto centralizado (estilo slide 2 do REC)."""
    txBox = slide.shapes.add_textbox(Cm(7), Cm(5.0), Cm(11.4), Cm(4.05))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = text
    run.font.size = Pt(32); run.font.bold = True; run.font.color.rgb = BLUE


def _fit_image(path, max_w, max_h):
    with PILImage.open(path) as im:
        iw, ih = im.size
    scale = min(max_w / (iw/150*2.54), max_h / (ih/150*2.54), 1.0)
    return iw/150*2.54*scale, ih/150*2.54*scale


def add_full_image(slide, path, top=1.4):
    """Imagem full-width abaixo do título (estilo slide 7/9 do REC)."""
    max_w = SW - 0.84
    max_h = SH - top - 0.3
    fw, fh = _fit_image(path, max_w, max_h)
    left = (SW - fw) / 2
    vtop = top + (max_h - fh) / 2
    slide.shapes.add_picture(path, Cm(left), Cm(vtop), Cm(fw), Cm(fh))


def add_two_images(slide, p1, p2, top=1.4, gap=0.5):
    """Duas imagens lado a lado (estilo slide 12/14 do REC)."""
    each_w = (SW - 0.84 - gap) / 2
    max_h  = SH - top - 0.3
    for path, x0 in [(p1, 0.42), (p2, 0.42 + each_w + gap)]:
        if not os.path.exists(path):
            print(f'  !! não encontrado: {path}'); continue
        fw, fh = _fit_image(path, each_w, max_h)
        vtop = top + (max_h - fh) / 2
        slide.shapes.add_picture(path, Cm(x0), Cm(vtop), Cm(fw), Cm(fh))


def add_text_below_title(slide, lines, top=1.4, height=2.5):
    """Bloco de texto abaixo do título (estilo slide 3 do REC)."""
    txBox = slide.shapes.add_textbox(Cm(1.5), Cm(top), Cm(SW - 3), Cm(height))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        is_header = line.endswith(':') or line.isupper()
        run.text = line
        run.font.size = Pt(11)
        run.font.bold = is_header
        run.font.color.rgb = BLUE if is_header else DARK


# ── Abre PPTX base (REC PADRÕES) ────────────────────────────
prs    = Presentation(ORIG)
slides = prs.slides

# ── SLIDE 1 — Título ─────────────────────────────────────────
s = slides[0]; clear_slide(s)
add_center_title(
    s,
    'Random Forest',
    'Detecção de Área Queimada — Imagens Landsat 8/9',
    'Ryan Alves · Matheus Sales · Pedro · Abner'
)
print('Slide 1: Título')

# ── SLIDE 2 — Seção ──────────────────────────────────────────
s = slides[1]; clear_slide(s)
add_section_title(s, 'Pipeline')
print('Slide 2: Seção')

# ── SLIDE 3 — Como Funciona ───────────────────────────────────
s = slides[2]; clear_slide(s)
add_title(s, 'Random Forest')
add_text_below_title(s, [
    'O Random Forest combina centenas de árvores de decisão:',
    '• Cada árvore aprende com uma amostra de pixels e atributos aleatórios',
    '• A baixa correlação entre árvores garante que os erros se cancelem',
    '• Saída: probabilidade de queimado (0–1) para cada pixel',
], top=1.4, height=2.6)
add_full_image(s, f'{GEN}/fig_rf_diversidade.png', top=4.2)
print('Slide 3: Como Funciona')

# ── SLIDE 4 — Entradas ────────────────────────────────────────
s = slides[3]; clear_slide(s)
add_title(s, 'Entradas')
add_full_image(s, PIPELINE_IMG, top=1.4)
print('Slide 4: Entradas — pipeline')

# ── SLIDE 5 — Features (separação espectral) ─────────────────
s = slides[4]; clear_slide(s)
add_title(s, 'Entradas')
add_full_image(s, f'{GEN}/fig_features_sep.png', top=1.4)
print('Slide 5: Features')

# ── SLIDE 6 — Dados Desbalanceados ───────────────────────────
s = slides[5]; clear_slide(s)
add_title(s, 'Dados desbalanceados')
add_text_below_title(s, [
    '4–8% dos pixels são queimados → modelo aprende a ignorar a minoria',
    '4 estratégias testadas: class_weight · Undersampling · SMOTE · SMOTETomek',
], top=1.4, height=1.5)
add_full_image(s, f'{GEN}/fig_resampling.png', top=3.1)
print('Slide 6: Desbalanceamento')

# ── SLIDE 7 — Validação Espacial ─────────────────────────────
s = slides[6]; clear_slide(s)
add_title(s, 'Validação por Blocos Espaciais')
add_text_below_title(s, [
    'Pixels vizinhos são altamente correlacionados — divisão aleatória infla métricas',
    'Solução: grade 4×4 (16 blocos); 4 blocos (~25%) reservados para teste',
], top=1.4, height=1.5)
add_full_image(s, f'{RF_FIGS}/fig_blocos_espaciais.png', top=3.1)
print('Slide 7: Blocos Espaciais')

# ── SLIDE 8 — Estratégias — Comparação ───────────────────────
s = slides[7]; clear_slide(s)
add_title(s, 'Estratégias de Desbalanceamento')
add_two_images(s, f'{RF_FIGS}/fig_estrategias_comp.png', TAB_ESTRATEGIAS,
               top=1.4)
print('Slide 8: Estratégias')

# ── SLIDE 9 — Threshold ───────────────────────────────────────
s = slides[8]; clear_slide(s)
add_title(s, 'Otimização do Limiar de Decisão')
add_two_images(s, f'{RF_FIGS}/fig_threshold_pr.png', TAB_THRESHOLD,
               top=1.4)
print('Slide 9: Threshold')

# ── SLIDE 10 — Validação Cross-cena ──────────────────────────
s = slides[9]; clear_slide(s)
add_title(s, 'Validação Cross-Cena')
add_text_below_title(s, [
    'Treino: cena 221/074  →  Teste: cena 220/075 (nunca vista)',
    'F₁ = 0,411  ·  Precision = 48,6%  ·  MCC = 0,374',
], top=1.4, height=1.5)
add_two_images(s, f'{RF_FIGS}/fig_cena_221074.png',
                  f'{RF_FIGS}/fig_cena_220075.png', top=3.1)
print('Slide 10: Cross-cena')

# ── SLIDE 11 — Saída: Mapa de Predição ───────────────────────
s = slides[10]; clear_slide(s)
add_title(s, 'Saída')
add_full_image(s, f'{RF_FIGS}/fig_mapa_predicao.png', top=1.4)
print('Slide 11: Mapa de predição')

# ── SLIDE 12 — Saída: Métricas + Mapa erros ──────────────────
s = slides[11]; clear_slide(s)
add_title(s, 'Saída')
add_two_images(s, f'{RF_FIGS}/fig_mapa_erros.png', TAB_MAPA, top=1.4)
print('Slide 12: Erros + métricas')

# ── SLIDE 13 — Importância de Atributos ──────────────────────
s = slides[12]; clear_slide(s)
add_title(s, 'Importância dos Atributos')
add_text_below_title(s, [
    'Distribuição uniforme (8,3%–11,5%) confirma ausência de leakage de dados',
    'Sem dNBR: nenhum atributo domina → modelo aprende de verdade',
], top=1.4, height=1.5)
add_full_image(s, f'{RF_FIGS}/fig_importancia.png', top=3.1)
print('Slide 13: Importância')

# ── SLIDES 14-18 — Limpos ────────────────────────────────────
for i in range(13, len(slides)):
    clear_slide(slides[i])
print('Slides 14-18: limpos')

# ── Salvar ────────────────────────────────────────────────────
prs.save(OUT)
print(f'\n✓ Salvo: {OUT}')
print(f'  Slides com conteúdo: 13 / {len(prs.slides)}')
