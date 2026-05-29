"""
PPTX com explicações escritas + gráficos.
Layout: título no topo | imagem à esquerda (60%) | texto à direita (37%)
"""
import os, warnings
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
warnings.filterwarnings('ignore')

from sklearn.datasets import make_classification
from sklearn.tree import DecisionTreeClassifier, plot_tree
from sklearn.ensemble import RandomForestClassifier, BaggingClassifier
from sklearn.model_selection import train_test_split
from sklearn.metrics import (fbeta_score, matthews_corrcoef, cohen_kappa_score,
                             precision_recall_curve, average_precision_score)
from sklearn.metrics import recall_score, precision_score, f1_score
from imblearn.over_sampling import SMOTE
from imblearn.combine import SMOTETomek
from imblearn.under_sampling import RandomUnderSampler
from scipy.ndimage import gaussian_filter

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image as PILImage

np.random.seed(42)

FIGS    = "/home/matheus-sales/Documents/atividade_reconhecimento/relatorio/pptx_figs"
RF_FIGS = "/home/matheus-sales/Documents/atividade_reconhecimento/relatorio/figs/rf"
ORIG    = "/home/matheus-sales/Documents/atividade_reconhecimento/relatorio/Classificação de áreas queimadas - REC PADRÕES.pptx"
OUT     = "/home/matheus-sales/Documents/atividade_reconhecimento/relatorio/RF_apresentacao.pptx"
os.makedirs(FIGS, exist_ok=True)

# ══════════════════════════════════════════════════════════
# PARTE 1 — FIGURAS
# ══════════════════════════════════════════════════════════
def savefig(name):
    path = f"{FIGS}/{name}.png"
    plt.savefig(path, dpi=150, bbox_inches='tight')
    plt.close('all')
    print(f"  ✓ {name}.png")

print("=== Gerando figuras ===")

# FIG 1: RF vs Bagging
X_s,y_s = make_classification(n_samples=5000,n_features=12,n_informative=5,
    n_redundant=3,n_classes=2,weights=[0.95,0.05],random_state=42)
Xtr,Xte,ytr,yte = train_test_split(X_s,y_s,test_size=0.3,stratify=y_s,random_state=42)
bag = BaggingClassifier(estimator=DecisionTreeClassifier(),n_estimators=100,random_state=42).fit(Xtr,ytr)
rff = RandomForestClassifier(n_estimators=100,random_state=42).fit(Xtr,ytr)
pb = np.array([t.predict(Xte) for t in bag.estimators_])
pr = np.array([t.predict(Xte) for t in rff.estimators_])
cb,cr = np.corrcoef(pb),np.corrcoef(pr)
fig,axes = plt.subplots(1,2,figsize=(11,4))
for ax,mat,tit in [(axes[0],cb,'Bagging'),(axes[1],cr,'Random Forest')]:
    im=ax.imshow(mat,cmap='hot_r',vmin=0,vmax=1)
    ax.set_title(f'{tit}\nCorr. média={np.triu(mat,1).mean():.3f}',fontsize=11)
    ax.set_xlabel('Árvore'); ax.set_ylabel('Árvore')
    plt.colorbar(im,ax=ax,fraction=0.046)
plt.tight_layout()
savefig("fig_rf_diversidade")

# FIG 2: CART
n0,n1=5000,300
X_d=np.column_stack([np.concatenate([np.random.normal(0.05,0.15,n0),np.random.normal(0.75,0.20,n1)]),
                     np.concatenate([np.random.normal(0.40,0.15,n0),np.random.normal(-0.10,0.20,n1)])])
y_d=np.array([0]*n0+[1]*n1)
cart=DecisionTreeClassifier(max_depth=3,criterion='entropy',random_state=42).fit(X_d,y_d)
fig,axes=plt.subplots(1,2,figsize=(11,4))
h=0.01; xx,yy=np.meshgrid(np.arange(-0.5,1.5,h),np.arange(-0.8,1.0,h))
Z=cart.predict(np.c_[xx.ravel(),yy.ravel()]).reshape(xx.shape)
axes[0].contourf(xx,yy,Z,alpha=0.3,cmap=plt.cm.RdYlGn_r)
axes[0].scatter(X_d[y_d==0,0],X_d[y_d==0,1],s=2,alpha=0.3,c='green',label='Não queimado')
axes[0].scatter(X_d[y_d==1,0],X_d[y_d==1,1],s=8,alpha=0.6,c='red',label='Queimado')
axes[0].set_xlabel('dNBR'); axes[0].set_ylabel('NBR pré')
axes[0].set_title('Superfície de Decisão CART'); axes[0].axvline(0.65,color='k',ls='--',lw=1.5,label='0.65')
axes[0].legend(fontsize=8)
plot_tree(cart,feature_names=['dNBR','NBR_pre'],class_names=['Não q.','Queimado'],
          filled=True,rounded=True,ax=axes[1],fontsize=7,impurity=True)
axes[1].set_title('Estrutura da Árvore (prof.=3)')
plt.tight_layout()
savefig("fig_cart")

# FIG 3: Features
n0,n1=2000,120
fs={'pre_B4':(np.random.normal(800,150,n0),np.random.normal(600,200,n1)),
    'pre_B5':(np.random.normal(2200,400,n0),np.random.normal(800,300,n1)),
    'pre_B7':(np.random.normal(1200,200,n0),np.random.normal(1800,400,n1)),
    'pre_NDVI':(np.random.normal(0.45,0.15,n0),np.random.normal(0.25,0.15,n1)),
    'pre_NBR':(np.random.normal(0.35,0.15,n0),np.random.normal(0.10,0.20,n1)),
    'pos_B5':(np.random.normal(2100,400,n0),np.random.normal(600,250,n1)),
    'pos_NBR':(np.random.normal(0.33,0.15,n0),np.random.normal(-0.35,0.25,n1)),
    'dNBR':(np.random.normal(0.02,0.10,n0),np.random.normal(0.75,0.20,n1)),
    'dNDVI':(np.random.normal(0.01,0.08,n0),np.random.normal(0.45,0.15,n1))}
fig,axes=plt.subplots(3,3,figsize=(12,8))
axes=axes.ravel()
for ax,(nm,(v0,v1)) in zip(axes,fs.items()):
    bins=np.linspace(min(v0.min(),v1.min()),max(v0.max(),v1.max()),60)
    ax.hist(v0,bins=bins,alpha=0.6,color='#2ecc71',density=True,label='Não queimado')
    ax.hist(v1,bins=bins,alpha=0.7,color='#e74c3c',density=True,label='Queimado')
    sep=abs(v0.mean()-v1.mean())/(np.sqrt((v0.std()**2+v1.std()**2)/2)+1e-6)
    ax.set_title(f'{nm} (sep={sep:.1f})',fontsize=8); ax.tick_params(labelsize=7)
    if nm=='pre_B4': ax.legend(fontsize=7)
plt.tight_layout()
savefig("fig_features_sep")

# FIG 4: Reamostagem
Xmaj=np.random.multivariate_normal([0,0],[[1,0.3],[0.3,1]],500)
Xmin=np.random.multivariate_normal([2,2],[[0.4,0],[0,0.4]],30)
X2d=np.vstack([Xmaj,Xmin]); y2d=np.array([0]*500+[1]*30)
Xu,yu=RandomUnderSampler(random_state=42).fit_resample(X2d,y2d)
Xs,ys=SMOTE(random_state=42,k_neighbors=5).fit_resample(X2d,y2d)
Xt,yt=SMOTETomek(random_state=42).fit_resample(X2d,y2d)
fig,axes=plt.subplots(2,2,figsize=(10,8))
for ax,Xp,yp,tit in [
    (axes[0,0],X2d,y2d,f'Original\n500 não-q. / 30 queimado'),
    (axes[0,1],Xu,yu,f'Undersampling\n{(yu==0).sum()} / {(yu==1).sum()}'),
    (axes[1,0],Xs,ys,f'SMOTE\n{(ys==0).sum()} / {(ys==1).sum()}'),
    (axes[1,1],Xt,yt,f'SMOTETomek\n{(yt==0).sum()} / {(yt==1).sum()}')]:
    ax.scatter(Xp[yp==0,0],Xp[yp==0,1],c='#2ecc71',s=12,alpha=0.4)
    ax.scatter(Xp[yp==1,0],Xp[yp==1,1],c='#e74c3c',s=30,alpha=0.8)
    ax.set_title(tit,fontsize=9); ax.set_xlim(-4,5); ax.set_ylim(-4,5); ax.grid(alpha=0.2)
plt.tight_layout()
savefig("fig_resampling")

# FIG 5: Blocos espaciais
ruido=np.random.randn(100,100)
campo=gaussian_filter(ruido,sigma=10)
campo=(campo-campo.min())/(campo.max()-campo.min())
label_img=(campo>0.65).astype(int)
Ximg=campo.ravel().reshape(-1,1)+np.random.randn(10000,1)*0.05; yimg=label_img.ravel()
coords=np.argwhere(np.ones((100,100),bool))
idx_al_tr=np.random.choice(len(yimg),int(0.7*len(yimg)),replace=False)
idx_al_te=np.setdiff1d(np.arange(len(yimg)),idx_al_tr)
bl=(coords[:,0]//25).clip(0,3)*4+(coords[:,1]//25).clip(0,3)
idx_bl_te=np.where(np.isin(bl,[0,5,10,15]))[0]
idx_bl_tr=np.setdiff1d(np.arange(len(yimg)),idx_bl_te)
res={}
for nm,itr,ite in [('Aleatória',idx_al_tr,idx_al_te),('Blocos',idx_bl_tr,idx_bl_te)]:
    rfs=RandomForestClassifier(n_estimators=50,random_state=42).fit(Ximg[itr],yimg[itr])
    res[nm]=fbeta_score(yimg[ite],rfs.predict(Ximg[ite]),beta=2,zero_division=0)
fig,axes=plt.subplots(1,3,figsize=(12,4))
axes[0].imshow(campo,cmap='YlGn'); axes[0].contour(label_img,levels=[0.5],colors='red',lw=2)
axes[0].set_title('Imagem simulada\n(vermelho=queimado)'); axes[0].axis('off')
ma=np.zeros(10000); ma[idx_al_tr]=0.3; ma[idx_al_te]=0.8
axes[1].imshow(ma.reshape(100,100),cmap='RdBu',vmin=0,vmax=1)
axes[1].set_title(f'Divisão Aleatória\nF2={res["Aleatória"]:.3f} (inflado!)'); axes[1].axis('off')
mb=np.zeros(10000); mb[idx_bl_tr]=0.3; mb[idx_bl_te]=0.8
axes[2].imshow(mb.reshape(100,100),cmap='RdBu',vmin=0,vmax=1)
axes[2].set_title(f'Blocos Espaciais\nF2={res["Blocos"]:.3f} (realista)'); axes[2].axis('off')
plt.tight_layout()
savefig("fig_blocos")

# FIG 6: Threshold
Xt2,yt2=make_classification(n_samples=5000,n_features=8,n_informative=4,
    n_classes=2,weights=[0.95,0.05],random_state=42)
Xtr2,Xte2,ytr2,yte2=train_test_split(Xt2,yt2,test_size=0.3,stratify=yt2,random_state=42)
rft=RandomForestClassifier(n_estimators=100,class_weight='balanced',random_state=42).fit(Xtr2,ytr2)
ypr=rft.predict_proba(Xte2)[:,1]
prec,rec,thr=precision_recall_curve(yte2,ypr)
df2=4*prec+rec; f2c=np.where(df2>0,5*prec*rec/df2,0)
thr_opt=thr[np.argmax(f2c[:-1])]
i05=np.argmin(np.abs(thr-0.5)); if2=np.argmax(f2c[:-1])
fig,axes=plt.subplots(1,2,figsize=(11,4))
axes[0].plot(rec,prec,'b-',lw=2,label=f'AUPRC={average_precision_score(yte2,ypr):.3f}')
axes[0].scatter(rec[i05],prec[i05],c='gray',s=100,zorder=5,label='τ=0.50')
axes[0].scatter(rec[if2],prec[if2],c='red',s=100,zorder=5,label=f'τ*={thr_opt:.2f}')
axes[0].set_xlabel('Recall'); axes[0].set_ylabel('Precision')
axes[0].set_title('Curva Precision-Recall'); axes[0].legend(fontsize=9); axes[0].grid(alpha=0.3)
axes[1].plot(thr,f2c[:-1],'red',lw=2)
axes[1].axvline(0.5,color='gray',ls='--',label='τ=0.50')
axes[1].axvline(thr_opt,color='red',ls=':',label=f'τ*={thr_opt:.2f}')
axes[1].set_xlabel('Threshold'); axes[1].set_ylabel('F2-Score')
axes[1].set_title('F2-Score por Threshold'); axes[1].legend(fontsize=9); axes[1].grid(alpha=0.3)
plt.tight_layout()
savefig("fig_threshold")

# FIG 7: Métricas
TN,FP,FN,TP=9650,50,75,225
fig,axes=plt.subplots(1,2,figsize=(11,4))
axes[0].imshow([[TN,FP],[FN,TP]],cmap='Blues')
axes[0].set_xticks([0,1]); axes[0].set_xticklabels(['Pred: Não q.','Pred: Queimado'])
axes[0].set_yticks([0,1]); axes[0].set_yticklabels(['Real: Não q.','Real: Queimado'])
axes[0].text(0,0,f'TN={TN}',ha='center',va='center',fontsize=11,color='navy')
axes[0].text(1,0,f'FP={FP}\nalarme',ha='center',va='center',fontsize=11,color='darkred')
axes[0].text(0,1,f'FN={FN}\nnão detect.',ha='center',va='center',fontsize=11,color='darkred')
axes[0].text(1,1,f'TP={TP}',ha='center',va='center',fontsize=11,color='darkgreen')
axes[0].set_title('Matriz de Confusão (97:3)',fontsize=10)
mets={'OA':(TN+TP)/(TN+FP+FN+TP),'Precision':TP/(TP+FP),'Recall':TP/(TP+FN),
      'F1':2*TP/(2*TP+FP+FN),'F2':5*TP/(5*TP+4*FN+FP),
      'Kappa':cohen_kappa_score([0]*9700+[1]*300,[0]*9650+[1]*50+[0]*75+[1]*225),
      'MCC':(TP*TN-FP*FN)/np.sqrt((TP+FP)*(TP+FN)*(TN+FP)*(TN+FN))}
cors=['#e74c3c','#3498db','#2ecc71','#f39c12','#e67e22','#9b59b6','#1abc9c']
axes[1].barh(list(mets.keys()),list(mets.values()),color=cors,alpha=0.85)
for i,(nm,v) in enumerate(mets.items()): axes[1].text(v+0.01,i,f'{v:.3f}',va='center',fontsize=9)
axes[1].set_xlim(0,1.15); axes[1].set_xlabel('Valor')
axes[1].set_title('Métricas com desbalanceamento 97:3',fontsize=10)
axes[1].axvline(1.0,color='gray',ls='--',lw=1); axes[1].grid(axis='x',alpha=0.3)
plt.tight_layout()
savefig("fig_metricas")

# FIG 8: Importância técnica
fnames=['pre_b4','pre_b5','pre_b7','pre_ndvi','pre_nbr',
        'pos_b4','pos_b5','pos_b7','pos_ndvi','pos_nbr','dnbr','dndvi']
n0,n1=3000,180
Xf=np.random.randn(n0+n1,12)
Xf[:n0,10]=np.random.normal(0.02,0.10,n0); Xf[n0:,10]=np.random.normal(0.80,0.18,n1)
Xf[:n0,9]=np.random.normal(0.35,0.12,n0);  Xf[n0:,9]=np.random.normal(-0.25,0.20,n1)
Xf[:n0,4]=np.random.normal(0.40,0.15,n0);  Xf[n0:,4]=np.random.normal(0.20,0.18,n1)
yf=np.array([0]*n0+[1]*n1)
Xft,Xfe,yft,yfe=train_test_split(Xf,yf,test_size=0.3,stratify=yf,random_state=42)
rfi=RandomForestClassifier(n_estimators=200,class_weight='balanced',random_state=42).fit(Xft,yft)
imp=rfi.feature_importances_; iord=np.argsort(imp)[::-1]
def corf(nm): return '#e74c3c' if nm.startswith('d') else ('#3498db' if 'pre' in nm else '#f39c12')
fig,axes=plt.subplots(1,2,figsize=(11,4))
axes[0].bar(range(12),imp[iord],color=[corf(fnames[i]) for i in iord],alpha=0.85)
axes[0].set_xticks(range(12)); axes[0].set_xticklabels([fnames[i] for i in iord],rotation=45,ha='right',fontsize=8)
axes[0].set_ylabel('Importância (Gini)'); axes[0].set_title('Importância por Feature')
axes[0].grid(axis='y',alpha=0.3)
axes[0].legend(handles=[mpatches.Patch(color='#e74c3c',label='Δ temporal'),
    mpatches.Patch(color='#3498db',label='Pré'),mpatches.Patch(color='#f39c12',label='Pós')],fontsize=8)
ia=np.cumsum(imp[iord])
axes[1].plot(range(1,13),ia,'b-o',ms=5)
axes[1].axhline(0.90,color='red',ls='--',label='90%'); axes[1].axhline(0.95,color='orange',ls='--',label='95%')
n90=np.argmax(ia>=0.90)+1; axes[1].axvline(n90,color='red',ls=':',label=f'{n90} feat→90%')
axes[1].set_xlabel('Nº features'); axes[1].set_ylabel('Acumulada')
axes[1].set_title('Importância Acumulada'); axes[1].legend(fontsize=8); axes[1].set_xticks(range(1,13)); axes[1].grid(alpha=0.3)
plt.tight_layout()
savefig("fig_importancia_tecnica")

# FIG 9: Hyperopt
Ne=np.linspace(50,400,100); Ml=np.linspace(1,20,100)
N,M=np.meshgrid(Ne,Ml)
perf=(np.exp(-((N-200)**2)/(2*80**2))*np.exp(-((M-5)**2)/(2*3**2))*0.4+0.5
      +np.random.randn(*N.shape)*0.02)
gs_pts=[(n,m) for n in [100,200,300] for m in [1,5,10,15]]
rs_n=np.random.uniform(50,400,20); rs_m=np.random.uniform(1,20,20)
fig,axes=plt.subplots(1,2,figsize=(11,4))
for ax,tit,px,py in [
    (axes[0],'Grid Search (12 fixos)',[p[0] for p in gs_pts],[p[1] for p in gs_pts]),
    (axes[1],'Randomized Search (20 aleat.)',rs_n,rs_m)]:
    im=ax.contourf(N,M,perf,levels=20,cmap='viridis')
    ax.scatter(px,py,c='white',s=60,zorder=5,edgecolors='red',lw=1.5,label='Avaliados')
    ax.set_xlabel('n_estimators'); ax.set_ylabel('min_samples_leaf')
    ax.set_title(tit,fontsize=10); ax.legend(fontsize=8)
    plt.colorbar(im,ax=ax,fraction=0.046,label='F2')
plt.tight_layout()
savefig("fig_hyperopt")

# FIG 10: Estratégias simulado
Xb,yb=make_classification(n_samples=5000,n_features=12,n_informative=5,n_redundant=3,
    n_classes=2,weights=[0.95,0.05],random_state=42)
Xtr3,Xte3,ytr3,yte3=train_test_split(Xb,yb,test_size=0.3,stratify=yb,random_state=42)
ests={'Sem tratamento':(Xtr3,ytr3,{}),
      'class_weight':(Xtr3,ytr3,{'class_weight':'balanced'}),
      'Undersampling':(*RandomUnderSampler(random_state=42).fit_resample(Xtr3,ytr3),{}),
      'SMOTE':(*SMOTE(random_state=42).fit_resample(Xtr3,ytr3),{}),
      'SMOTETomek':(*SMOTETomek(random_state=42).fit_resample(Xtr3,ytr3),{})}
rs2={'F2':[],'MCC':[],'Recall':[],'Precision':[],'label':[]}
for nm,(Xr,yr,kw) in ests.items():
    rfs=RandomForestClassifier(n_estimators=100,random_state=42,n_jobs=-1,**kw).fit(Xr,yr)
    yp=rfs.predict(Xte3)
    rs2['F2'].append(fbeta_score(yte3,yp,beta=2,zero_division=0))
    rs2['MCC'].append(matthews_corrcoef(yte3,yp))
    rs2['Recall'].append(recall_score(yte3,yp,pos_label=1,zero_division=0))
    rs2['Precision'].append(precision_score(yte3,yp,pos_label=1,zero_division=0))
    rs2['label'].append(nm)
mpl2=['F2','MCC','Recall','Precision']
x=np.arange(len(mpl2)); w=0.15; cs=['#95a5a6','#3498db','#e67e22','#2ecc71','#9b59b6']
fig,ax=plt.subplots(figsize=(11,4))
for i,lbl in enumerate(rs2['label']):
    ax.bar(x+i*w,[rs2[m][i] for m in mpl2],w,label=lbl,color=cs[i],alpha=0.85)
ax.set_xticks(x+w*2); ax.set_xticklabels(mpl2,fontsize=11)
ax.set_ylabel('Valor'); ax.set_ylim(0,1.05)
ax.set_title('Comparação de Estratégias (dados simulados 19:1)')
ax.legend(fontsize=8); ax.grid(axis='y',alpha=0.3)
plt.tight_layout()
savefig("fig_estrategias_sim")

print("✓ Figuras OK\n")


# ══════════════════════════════════════════════════════════
# PARTE 2 — CONSTRUIR PPTX
# ══════════════════════════════════════════════════════════
print("=== Construindo PPTX ===")

P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'

# Cores
BLUE  = RGBColor(0x42, 0x85, 0xF4)   # Google Blue do tema
DARK  = RGBColor(0x21, 0x21, 0x21)   # quase preto
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY  = RGBColor(0x55, 0x55, 0x55)

# Layout: imagem à esquerda, texto à direita
IMG_LEFT  = 0.4    # cm
IMG_TOP   = 1.40   # cm (abaixo do título)
IMG_W     = 14.8   # cm (58% do slide)
IMG_H     = 12.5   # cm

TXT_LEFT  = 15.6   # cm
TXT_TOP   = 1.40   # cm
TXT_W     = 9.4    # cm
TXT_H     = 12.5   # cm

SW, SH = 25.40, 14.29   # slide width/height em cm


def clear_slide_content(slide):
    sp_tree = slide.shapes._spTree
    to_remove = []
    for child in list(sp_tree):
        tag = etree.QName(child.tag).localname
        if tag == 'pic':
            to_remove.append(child)
        elif tag == 'sp':
            nvPr = child.find(f'.//{{{P_NS}}}nvPr')
            ph = nvPr.find(f'{{{P_NS}}}ph') if nvPr is not None else None
            if ph is None:
                to_remove.append(child)
        elif tag in ('grpSp','graphicFrame','cxnSp'):
            to_remove.append(child)
    for c in to_remove:
        sp_tree.remove(c)
    for ph in slide.placeholders:
        if ph.has_text_frame:
            ph.text_frame.clear()


def set_title(slide, text, fontsize=16):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx in (0, 1):
            tf = ph.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            run = p.add_run(); run.text = text
            run.font.size = Pt(fontsize); run.font.bold = True
            run.font.color.rgb = BLACK
            return


def set_section(slide, text, fontsize=30):
    for ph in slide.placeholders:
        tf = ph.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = text
        run.font.size = Pt(fontsize); run.font.bold = True
        run.font.color.rgb = BLACK
        return


def add_img_left(slide, path, top=IMG_TOP, left=IMG_LEFT, max_w=IMG_W, max_h=IMG_H):
    if not os.path.exists(path):
        print(f"  !! Não encontrado: {path}"); return
    with PILImage.open(path) as im:
        iw, ih = im.size
    img_w, img_h = iw/150*2.54, ih/150*2.54
    scale = min(max_w/img_w, max_h/img_h, 1.0)
    fw, fh = img_w*scale, img_h*scale
    # centra verticalmente
    vtop = top + (max_h - fh) / 2
    slide.shapes.add_picture(path, Cm(left), Cm(vtop), Cm(fw), Cm(fh))


def add_text_panel(slide, title, bullets,
                   left=TXT_LEFT, top=TXT_TOP, width=TXT_W, height=TXT_H):
    txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = txBox.text_frame; tf.word_wrap = True

    # Título do painel
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(4)
    run = p.add_run(); run.text = title
    run.font.size = Pt(11); run.font.bold = True
    run.font.color.rgb = BLUE

    # Linha separadora (horizontal rule via texto vazio com border — simulada por espaçamento)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(2); p2.space_after = Pt(6)
    run2 = p2.add_run(); run2.text = '─' * 32
    run2.font.size = Pt(7); run2.font.color.rgb = RGBColor(0xA4, 0xC2, 0xF4)

    # Bullets
    for bullet in bullets:
        pb = tf.add_paragraph(); pb.alignment = PP_ALIGN.LEFT
        pb.space_before = Pt(3); pb.space_after = Pt(1)
        # Indent level
        pb.level = 0
        if bullet.startswith('  '):
            pb.level = 1
            bullet = bullet.strip()
        run = pb.add_run()
        # Bullet character
        run.text = ('   • ' if pb.level == 0 else '      ◦ ') + bullet
        run.font.size = Pt(9.5)
        run.font.color.rgb = DARK


def add_two_imgs_centered(slide, p1, p2, top=IMG_TOP, gap=0.4):
    each_w = (SW - 0.8 - gap) / 2
    max_h  = SH - top - 0.4
    for path, x0 in [(p1, 0.4), (p2, 0.4 + each_w + gap)]:
        if not os.path.exists(path):
            print(f"  !! {path}"); continue
        with PILImage.open(path) as im:
            iw, ih = im.size
        img_w, img_h = iw/150*2.54, ih/150*2.54
        sc = min(each_w/img_w, max_h/img_h, 1.0)
        fw, fh = img_w*sc, img_h*sc
        vtop = top + (max_h-fh)/2
        slide.shapes.add_picture(path, Cm(x0), Cm(vtop), Cm(fw), Cm(fh))


# ── Abre original e modifica in-place ───────────────────
prs = Presentation(ORIG)
slides = prs.slides


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 1 — Título
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = slides[0]; clear_slide_content(s)
for ph in s.placeholders:
    idx = ph.placeholder_format.idx
    tf = ph.text_frame; tf.word_wrap = True
    if idx == 0:
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = 'Random Forest para Detecção de Área Queimada'
        run.font.size = Pt(26); run.font.bold = True
    elif idx == 1:
        for line in ['Classificação binária pixel a pixel: Queimado (1) / Não Queimado (0)',
                     'Imagens Landsat 8 e 9 — Bandas B4, B5, B7 + índices NDVI e NBR',
                     'Reconhecimento de Padrões — UNESP ICT',
                     'Ryan Alves · Matheus Sales · Pedro · Abner']:
            p = tf.add_paragraph() if ph.text_frame.text else tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run(); run.text = line; run.font.size = Pt(16)
print("Slide 1: Título")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 2 — Entrada: Dados e Atributos
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = slides[1]; clear_slide_content(s)
set_title(s, 'Entrada — Dados Landsat e Atributos por Pixel')
add_img_left(s, f"{FIGS}/fig_features_sep.png")
add_text_panel(s, 'Entrada: o que alimenta o modelo', [
    'Fonte: 2 imagens Landsat (antes e depois do incêndio)',
    '  Cena 221/074 — treino e validação (4,8% queimado)',
    '  Cena 220/075 — teste final em área nova (7,9%)',
    '10 atributos extraídos por pixel:',
    '  B4, B5, B7, NDVI, NBR da imagem pré-fogo',
    '  B4, B5, B7, NDVI, NBR da imagem pós-fogo',
    'O gráfico mostra separação entre classes por atributo',
    '  verde = não queimado / vermelho = queimado',
    'dNBR e dNDVI foram excluídos (vazamento de dados):',
    '  O gabarito foi criado usando dNBR > 0.65',
    '  Usar dNBR como entrada = dar a resposta ao modelo',
])
print("Slide 2: Entrada")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 3 — Processamento: O Modelo RF
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = slides[2]; clear_slide_content(s)
set_title(s, 'Processamento — O Modelo Random Forest')
add_img_left(s, f"{FIGS}/fig_rf_diversidade.png")
add_text_panel(s, 'Como o RF processa cada pixel', [
    'Entrada: vetor de 10 atributos por pixel',
    'Processamento:',
    '  171 árvores de decisão treinadas em paralelo',
    '  Cada árvore aprende com uma amostra dos pixels',
    '  Cada divisão testa apenas alguns atributos ao acaso',
    '  Árvores diferentes cometem erros diferentes',
    'Saída: probabilidade de queimado (0 a 1) por pixel',
    '  A votação das árvores reduz o erro individual',
    'O gráfico mostra a correlação entre árvores:',
    '  RF (direita): árvores mais independentes',
    '  Bagging (esquerda): árvores mais correlacionadas',
    '  Mais independência → erros se cancelam melhor',
])
print("Slide 3: RF")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 4 — Processamento: Desbalanceamento
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = slides[3]; clear_slide_content(s)
set_title(s, 'Processamento — Corrigindo o Desbalanceamento de Classes')
add_img_left(s, f"{FIGS}/fig_resampling.png")
add_text_panel(s, 'Problema: 92–96% dos pixels não são queimados', [
    'Sem correção: modelo diz "não queimado" sempre',
    '  95% de acurácia, zero queimadas detectadas',
    '4 estratégias testadas para corrigir:',
    '  class_weight: penaliza mais o erro na minoria',
    '  Undersampling: descarta não-queimados → treino 1:1',
    '  SMOTE: cria pixels queimados sintéticos',
    '  SMOTETomek: SMOTE + remove exemplos ambíguos',
    'O gráfico mostra o efeito de cada estratégia',
    '  sobre a distribuição no espaço de atributos',
    'Resultado: class_weight com limiar ajustado',
    '  liderou o F₁ final (0.262)',
])
print("Slide 4: Desbalanceamento")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 5 — Processamento: Validação + Limiar
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = slides[4]; clear_slide_content(s)
set_title(s, 'Processamento — Validação Espacial e Ajuste do Limiar')
add_img_left(s, f"{RF_FIGS}/fig_threshold_pr.png")
add_text_panel(s, 'Duas decisões críticas no processamento', [
    'Validação por blocos espaciais:',
    '  Pixels vizinhos são muito parecidos (mesma cena)',
    '  Divisão aleatória infla o resultado (F₁ ≈ 0.89)',
    '  Solução: 16 blocos geográficos, 4 só para teste',
    '  Resultado real com blocos: F₁ ≈ 0.48',
    'Ajuste do limiar de decisão τ:',
    '  RF produz probabilidade → τ define "queimado"',
    '  Limiar padrão 50% detecta quase nada',
    '  Buscamos τ* que maximiza F₁',
    'O gráfico mostra a curva Precision vs Recall:',
    '  τ* = 0.65 → F₁ = 0.329',
    '  Precision = 37,8% / Recall = 29,1%',
])
print("Slide 5: Validação + Limiar")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 6 — Resultado das Estratégias
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = slides[5]; clear_slide_content(s)
set_title(s, 'Comparação das Estratégias — Dados Reais')
add_img_left(s, f"{RF_FIGS}/fig_estrategias_comp.png")
add_text_panel(s, 'Entrada → Processamento → Saída por estratégia', [
    'Entrada: mesmos 10 atributos para todas',
    'Processamento: varia a estratégia de balanceamento',
    'Saída (F₁ com τ*):',
    '  Sem tratamento: F₁ < 0.03 — inaceitável',
    '  Undersampling: alto Recall, menor F₁ geral',
    '  SMOTE / SMOTETomek: resultado parecido entre os dois',
    '  class_weight + τ*: melhor F₁ (0.262)',
    'O ajuste do limiar τ* foi decisivo:',
    '  Sem ele, todas as estratégias ficam abaixo de 0.05',
    '  Com ele, class_weight alcançou 37,8% de Precision',
    'Conclusão: balancear os dados não é suficiente',
    '  o limiar precisa ser ajustado junto',
])
print("Slide 6: Estratégias")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 7 — Saída: Mapa de Predição
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = slides[6]; clear_slide_content(s)
set_title(s, 'Saída — Mapa de Predição do RF')
add_img_left(s, f"{RF_FIGS}/fig_mapa_predicao.png")
add_text_panel(s, 'Entrada → Processamento → Saída', [
    'Entrada: imagem completa (~39 milhões de pixels)',
    'Processamento: RF classifica cada pixel com τ* = 0.65',
    'Saída: mapa binário queimado / não queimado',
    'O mapa tem três partes (esq → dir):',
    '  dNBR: referência física da cicatriz de fogo',
    '  Gabarito: máscara criada com dNBR > 0.65',
    '  Predição RF: resultado do modelo',
    '    vermelho = pixel classificado como queimado',
    'F₁ = 0.247 no mapa completo',
    '  Precision = 29,2% / Recall = 21,3%',
    'Principais cicatrizes foram detectadas',
    '  Forma espacial das áreas preservada',
])
print("Slide 7: Mapa predição")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 8 — Saída: Erros e Conclusão
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = slides[7]; clear_slide_content(s)
set_title(s, 'Saída — Análise de Erros e Conclusão')
add_img_left(s, f"{RF_FIGS}/fig_mapa_erros.png")
add_text_panel(s, 'O que o mapa de erros revela', [
    'Verde: queimada detectada corretamente (TP = 16.252)',
    'Laranja: alarme falso (FP = 39.328)',
    '  Solo seco e pastagem têm reflexo parecido com queimada',
    'Vermelho: queimada perdida (FN = 59.905)',
    '  Cicatrizes fracas têm sinal parecido com vegetação',
    'Cross-scene (cena nunca vista pelo modelo):',
    '  F₁ = 0.411 — modelo generalizou para nova área',
    'Limitação central: sem dNBR a separação é difícil',
    '  dNBR foi excluído porque criou o gabarito',
    '  Esse é o custo de um modelo sem vazamento de dados',
    'Conclusão: RF detecta queimadas sem usar o dNBR',
    '  mas há tradeoff entre Precision e Recall',
])
print("Slide 8: Erros e conclusão")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDES 9-18 — Limpos (não usados)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
for i in range(8, len(slides)):
    clear_slide_content(slides[i])
print("Slides 9-18: limpos")


# ── Salvar ──────────────────────────────────────────────
prs.save(OUT)
print(f"\n✓ Salvo: {OUT}")
print(f"  Slides com conteúdo: 8 / {len(prs.slides)}")
